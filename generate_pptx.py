from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Helper function to create a title slide
def add_title_slide(title, subtitle, notes=""):
    slide_layout = prs.slide_layouts[0] # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle
    
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    return slide

# Helper function to create bullet point slides
def add_bullet_slide(title, points, notes=""):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_placeholder.text = title
    tf = body_shape.text_frame
    tf.text = points[0]
    
    for point in points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    return slide

# Slide 1
add_title_slide(
    "Perancangan dan Implementasi Sistem ERP Sederhana Berbasis Web sebagai Prototype",
    "Laporan Kerja Praktik\nOleh: Arjuna Gunatama Sihombing (2317051085)\nJurusan Ilmu Komputer, Universitas Lampung (2026)",
    "Selamat pagi/siang Bapak/Ibu dosen penguji/pembimbing. Perkenalkan, nama saya Arjuna Gunatama Sihombing. Pada kesempatan ini, saya akan mempresentasikan hasil kerja praktik saya yang berjudul 'Perancangan dan Implementasi Sistem ERP Sederhana Berbasis Web sebagai Prototype Pendukung Proses Bisnis Perusahaan'."
)

# Slide 2
add_bullet_slide(
    "Latar Belakang Masalah",
    [
        "Kebutuhan akan sistem terintegrasi (ERP) untuk mendukung proses bisnis.",
        "Pengelolaan data dan pemantauan kinerja masih terpencar atau manual.",
        "Kebutuhan dashboard monitoring terpusat untuk keputusan real-time.",
        "Solusi: Pengembangan prototype sistem ERP berbasis web."
    ],
    "Latar belakang dari proyek ini berawal dari kebutuhan perusahaan akan sistem yang terintegrasi. Saat ini, beberapa proses bisnis datanya masih tersebar sehingga sulit untuk dimonitoring secara menyeluruh. Oleh karena itu, saya merancang sebuah prototype Sistem ERP berbasis web yang dilengkapi dengan dashboard pemantauan kinerja."
)

# Slide 3
add_bullet_slide(
    "Tujuan & Manfaat",
    [
        "Tujuan: Membangun prototype ERP fungsional dan fitur dashboard.",
        "Tujuan: Mengaplikasikan metodologi RAD dalam dunia industri.",
        "Manfaat Perusahaan: Mendapatkan model pengelolaan data terpusat.",
        "Manfaat Mahasiswa: Pengalaman nyata dalam pengembangan full-stack."
    ],
    "Tujuan utama proyek ini adalah membangun purwarupa (prototype) sistem ERP yang berjalan dengan baik, khususnya pada fitur pemantauan (monitoring). Dengan adanya prototype ini, perusahaan mendapatkan gambaran nyata mengenai integrasi data, dan bagi saya pribadi, ini adalah pengalaman yang sangat berharga dalam menerapkan teori perancangan perangkat lunak."
)

# Slide 4
add_bullet_slide(
    "Metodologi Pengembangan (RAD)",
    [
        "Menggunakan Metode: Rapid Application Development (RAD).",
        "Karakteristik: Iteratif, Cepat, dan Fleksibel.",
        "Keterlibatan Pengguna: Menerima feedback secara langsung di setiap iterasi.",
        "Tahapan: Requirement Planning -> User Design -> Construction -> Cutover."
    ],
    "Dalam pengembangan sistem ini, saya menggunakan metode RAD atau Rapid Application Development. Alasan pemilihannya adalah karena metode ini sangat cocok untuk pengembangan prototype. RAD memungkinkan siklus kerja yang iteratif dan cepat, sehingga aplikasi berkembang tepat sasaran sesuai masukan."
)

# Slide 5
add_bullet_slide(
    "Analisis & Perancangan Sistem",
    [
        "Use Case Diagram: Memetakan pengguna (Admin, User) beserta hak akses.",
        "Activity Diagram (Swimlanes): Menggambarkan alur kerja rinci antar divisi.",
        "Entity Relationship Diagram (ERD): Mengatur struktur dan relasi database."
    ],
    "Sebelum tahap koding, sistem dirancang secara matang. Use Case Diagram untuk memetakan aktor. Activity Diagram dengan format swimlanes digunakan untuk alur kerja. Sedangkan ERD digunakan untuk memastikan struktur database seperti data karyawan sudah optimal."
)

# Slide 6
add_bullet_slide(
    "Implementasi Sistem",
    [
        "Teknologi Backend & Frontend: Framework Laravel (PHP).",
        "Teknologi Database: MySQL.",
        "Visualisasi Data: Library Chart.js / D3.js.",
        "Fitur Utama: Otentikasi aman, Role-Based Access Control (RBAC), Manajemen Data Karyawan."
    ],
    "Sistem ini dibangun menggunakan framework Laravel yang handal. Beberapa fitur utama yang berhasil dibuat antara lain sistem otentikasi dengan pembagian hak akses (RBAC), serta manajemen data karyawan."
)

# Slide 7
add_bullet_slide(
    "Highlight Fitur: Monitoring Dashboard",
    [
        "Visualisasi Data Real-time menggunakan grafik.",
        "Weekly Trends: Grafik garis untuk tren aktivitas mingguan.",
        "Status Doughnut: Persentase status pekerjaan berjalan.",
        "Top Performers: Peringkat karyawan dengan performa terbaik.",
        "Activity Heatmap: Menampilkan distribusi intensitas aktivitas."
    ],
    "Ini adalah fitur unggulan, yaitu Dashboard Pemantauan Administratif. Tujuannya merangkum data kompleks menjadi visual yang mudah dipahami. Mulai dari melihat tren mingguan hingga daftar peringkat performa karyawan. Sangat membantu manajer memonitor produktivitas secara real-time."
)

# Slide 8
add_bullet_slide(
    "Kesimpulan",
    [
        "Prototype ERP berhasil dibangun sesuai kebutuhan menggunakan metode RAD.",
        "Dokumentasi teknis (UML, ERD) berhasil diselesaikan dan mendukung sistem.",
        "Sistem mampu memvisualisasikan data kinerja secara dinamis untuk manajemen."
    ],
    "Sebagai kesimpulan, kerja praktik ini telah berhasil mencapai tujuannya, yaitu menghasilkan prototype sistem ERP yang fungsional. Penggunaan RAD terbukti efektif. Dashboard yang dibangun juga sudah bisa memvisualisasikan data."
)

# Slide 9
add_bullet_slide(
    "Saran & Pengembangan Lanjutan",
    [
        "Pengembangan Modul: Menambah modul Keuangan, HRD, dan Inventory.",
        "Peningkatan Infrastruktur: Optimasi database dan implementasi caching.",
        "Pengujian Lanjutan (UAT): Pengujian terhadap pengguna skala luas."
    ],
    "Saran pengembangan ke depannya adalah penambahan modul-modul lain sehingga sistem ERP ini menjadi lebih lengkap. Selain itu, diperlukan uji coba lanjutan (UAT) kepada seluruh staf terkait agar sistem benar-benar siap untuk tahap produksi."
)

# Slide 10
add_title_slide(
    "Terima Kasih",
    "Sesi Tanya Jawab (Q&A)",
    "Sekian presentasi dari saya terkait Laporan Kerja Praktik ini. Terima kasih atas perhatian Bapak/Ibu sekalian. Selanjutnya, saya kembalikan kepada moderator atau saya persilakan jika ada pertanyaan."
)

prs.save('Presentasi_KP_Arjuna.pptx')
print("PPTX generated successfully.")
